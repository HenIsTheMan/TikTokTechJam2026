#!/usr/bin/env python3
"""Generate a presentation for the Transformer GPU optimization benchmark."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path(__file__).with_name("transformer_optimization_pipeline_results.pptx")

# 16:9 slide and a restrained technical palette.
SW, SH = Inches(13.333), Inches(7.5)
BG = RGBColor(10, 18, 31)
PANEL = RGBColor(19, 31, 50)
PANEL_2 = RGBColor(26, 42, 65)
INK = RGBColor(238, 244, 251)
MUTED = RGBColor(158, 174, 194)
GRID = RGBColor(54, 73, 96)
TEAL = RGBColor(41, 211, 181)
CYAN = RGBColor(85, 178, 255)
AMBER = RGBColor(255, 190, 72)
CORAL = RGBColor(255, 111, 97)
WHITE = RGBColor(255, 255, 255)


def add_bg(slide, color=BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, size=18, color=INK, bold=False,
             font="Aptos", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
             margin=0.04, fit=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = Inches(margin)
    frame.margin_top = frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    # python-pptx font discovery is not available on every Linux build. Text
    # boxes are dimensioned explicitly, so automatic fit is unnecessary here.
    return box


def add_rich_text(slide, runs, x, y, w, h, size=18, align=PP_ALIGN.LEFT,
                  valign=MSO_ANCHOR.TOP, margin=0.04):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = Inches(margin)
    frame.margin_top = frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.alignment = align
    for text, color, bold in runs:
        run = p.add_run()
        run.text = text
        run.font.name = "Aptos"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_box(slide, x, y, w, h, fill=PANEL, line=GRID, radius=True, lw=1.0):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    box = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = line
    box.line.width = Pt(lw)
    return box


def add_label(slide, text, x, y, w, color=TEAL):
    box = add_box(slide, x, y, w, 0.32, fill=color, line=color)
    add_text(slide, text.upper(), x + 0.05, y + 0.015, w - 0.1, 0.27,
             size=10, color=BG, bold=True, align=PP_ALIGN.CENTER,
             valign=MSO_ANCHOR.MIDDLE, margin=0)
    return box


def add_title(slide, title, subtitle=None, section=None):
    if section:
        add_text(slide, section.upper(), 0.65, 0.35, 3.0, 0.25, size=10,
                 color=TEAL, bold=True, margin=0)
    add_text(slide, title, 0.65, 0.66, 12.0, 0.62, size=29, bold=True, margin=0)
    if subtitle:
        add_text(slide, subtitle, 0.65, 1.30, 12.0, 0.42, size=14,
                 color=MUTED, margin=0)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(1.82),
                                  Inches(12.03), Inches(0.018))
    line.fill.solid(); line.fill.fore_color.rgb = GRID; line.line.fill.background()


def add_footer(slide, number, source="Source: torch_transformer_benchmark.py"):
    add_text(slide, source, 0.65, 7.16, 10.8, 0.2, size=8.5, color=MUTED, margin=0)
    add_text(slide, f"{number:02d}", 12.1, 7.13, 0.55, 0.23, size=9,
             color=MUTED, bold=True, align=PP_ALIGN.RIGHT, margin=0)


def add_arrow(slide, x1, y1, x2, y2, color=TEAL, width=2.0):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    line.line.end_arrowhead = True
    return line


def add_bullets(slide, items, x, y, w, h, size=15, color=INK,
                bullet_color=TEAL, gap=8):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear(); tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.level = 0
        p.text = ""
        dot = p.add_run(); dot.text = "●  "; dot.font.color.rgb = bullet_color
        dot.font.size = Pt(size - 2)
        run = p.add_run(); run.text = item; run.font.name = "Aptos"
        run.font.size = Pt(size); run.font.color.rgb = color
    return box


def stage_box(slide, x, y, w, title, detail, color=CYAN, h=0.82):
    add_box(slide, x, y, w, h, fill=PANEL, line=color, lw=1.5)
    add_text(slide, title, x + 0.12, y + 0.10, w - 0.24, 0.25, size=13,
             color=color, bold=True, align=PP_ALIGN.CENTER, margin=0)
    add_text(slide, detail, x + 0.10, y + 0.42, w - 0.20, h - 0.48, size=10,
             color=MUTED, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)


def metric(slide, x, y, w, value, label, color=TEAL):
    add_box(slide, x, y, w, 1.05, fill=PANEL, line=GRID)
    add_text(slide, value, x + 0.1, y + 0.12, w - 0.2, 0.46, size=25,
             color=color, bold=True, align=PP_ALIGN.CENTER, margin=0)
    add_text(slide, label, x + 0.1, y + 0.66, w - 0.2, 0.22, size=10,
             color=MUTED, align=PP_ALIGN.CENTER, margin=0)


def add_slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
    # Decorative circuit-like tracks.
    for y, c, w in [(0.62, TEAL, 2.5), (0.88, CYAN, 1.65), (6.65, AMBER, 2.05)]:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.4), Inches(y),
                                      Inches(w), Inches(0.035))
        line.fill.solid(); line.fill.fore_color.rgb = c; line.line.fill.background()
    add_label(slide, "GPU optimization case study", 0.72, 0.72, 2.25, TEAL)
    add_text(slide, "Transformer inference,\noptimized stage by stage", 0.72, 1.38,
             8.7, 1.55, size=36, color=INK, bold=True, margin=0)
    add_text(slide, "Pipeline architecture · CUDA fusion · selective precision · measured results",
             0.74, 3.15, 9.2, 0.45, size=17, color=MUTED, margin=0)
    # Hero throughput path.
    add_box(slide, 0.74, 4.18, 11.9, 1.47, fill=PANEL, line=GRID)
    stages = [(1.02, "INPUT", "8 × 128 × 512", CYAN),
              (3.28, "6× BLOCK", "attention + FFN", TEAL),
              (5.64, "CUDA FUSION", "memory-bound epilogues", AMBER),
              (8.50, "HYBRID TF32", "safe FFN GEMMs", CORAL),
              (10.85, "OUTPUT", "normalized tokens", CYAN)]
    for i, (x, t, d, c) in enumerate(stages):
        stage_box(slide, x, 4.50, 1.55 if i != 2 else 2.05, t, d, c, h=0.80)
        if i < len(stages) - 1:
            next_x = stages[i + 1][0]
            add_arrow(slide, x + (1.55 if i != 2 else 2.05), 4.90, next_x - 0.12, 4.90, c)
    add_text(slide, "TikTok TechJam 2026 · PyTorch benchmark · target validation: RTX 5090",
             0.74, 6.73, 10.8, 0.27, size=10.5, color=MUTED, margin=0)
    add_text(slide, "01", 12.05, 6.70, 0.58, 0.28, size=10, color=MUTED,
             bold=True, align=PP_ALIGN.RIGHT, margin=0)


def add_slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
    add_title(slide, "The workload and the optimization target",
              "A fixed inference shape makes targeted kernel design practical.", "01 · Scope")
    metric(slide, 0.72, 2.15, 2.20, "8 × 128", "batch × sequence", CYAN)
    metric(slide, 3.10, 2.15, 2.20, "512", "model width", TEAL)
    metric(slide, 5.48, 2.15, 2.20, "8 × 64", "heads × head dim", AMBER)
    metric(slide, 7.86, 2.15, 2.20, "2048", "FFN width", CORAL)
    metric(slide, 10.24, 2.15, 2.20, "6", "Transformer layers", CYAN)
    add_box(slide, 0.72, 3.58, 5.86, 2.65, fill=PANEL, line=GRID)
    add_text(slide, "Optimization objective", 1.00, 3.88, 5.25, 0.34,
             size=18, bold=True, color=TEAL, margin=0)
    add_bullets(slide, [
        "Minimize end-to-end inference latency on CUDA.",
        "Preserve identical learned weights and output shape.",
        "Pass strict per-element numerical validation before benchmarking.",
    ], 1.00, 4.42, 5.25, 1.45, size=14)
    add_box(slide, 6.82, 3.58, 5.63, 2.65, fill=PANEL, line=GRID)
    add_text(slide, "Benchmark discipline", 7.10, 3.88, 5.10, 0.34,
             size=18, bold=True, color=CYAN, margin=0)
    add_bullets(slide, [
        "20 warmups; 100 repeats × 3 alternating rounds.",
        "CUDA Events time a fixed input on the current stream.",
        "Report median, mean, p90, minimum, throughput and peak allocation.",
    ], 7.10, 4.42, 5.02, 1.45, size=14, bullet_color=CYAN)
    add_footer(slide, 2, "Source: torch_transformer_benchmark.py defaults and benchmark_models()")


def add_slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
    add_title(slide, "End-to-end network pipeline",
              "Each block is pre-normalized; residual paths carry the token state through six layers.",
              "02 · Architecture")
    # Main pipeline.
    xs = [0.78, 2.18, 3.63, 5.22, 6.75, 8.26, 9.91, 11.35]
    stages = [
        ("X", "[B,S,512]", CYAN, 1.08), ("LN₁", "normalize", TEAL, 1.02),
        ("MHA", "Q·Kᵀ → softmax → V", AMBER, 1.28), ("+", "residual", CYAN, 0.90),
        ("LN₂", "normalize", TEAL, 1.02), ("FFN", "512→2048→512", CORAL, 1.26),
        ("+", "residual", CYAN, 0.90), ("NEXT", "repeat ×6", TEAL, 1.15),
    ]
    for i, ((title, detail, color, width), x) in enumerate(zip(stages, xs)):
        stage_box(slide, x, 2.45, width, title, detail, color, h=1.02)
        if i < len(stages) - 1:
            add_arrow(slide, x + width, 2.96, xs[i + 1] - 0.10, 2.96, color)
    # Residual arcs represented as clean lower lanes.
    add_text(slide, "residual lane", 1.02, 3.86, 1.05, 0.22, 9, MUTED, margin=0)
    add_arrow(slide, 1.12, 4.18, 5.65, 4.18, color=CYAN, width=1.5)
    add_text(slide, "residual lane", 5.92, 4.52, 1.05, 0.22, 9, MUTED, margin=0)
    add_arrow(slide, 5.85, 4.82, 10.25, 4.82, color=CYAN, width=1.5)
    add_box(slide, 0.78, 5.42, 11.72, 1.10, fill=PANEL, line=GRID)
    add_rich_text(slide, [
        ("After layer 6:  ", MUTED, False), ("Final LayerNorm", TEAL, True),
        ("  →  zero padded query rows (when masking is active)  →  ", MUTED, False),
        ("Y ∈ ℝ[B,S,512]", INK, True),
    ], 1.06, 5.74, 11.2, 0.44, size=16, align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE)
    add_footer(slide, 3, "Source: BaselineTransformerBlock and BaselineTransformer")


def add_slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
    add_title(slide, "Attention: collapse launches and avoid score materialization",
              "The optimized path delegates attention to PyTorch SDPA, which can dispatch to fused Flash Attention.",
              "03 · Stage optimizations")
    add_label(slide, "Baseline", 0.75, 2.08, 1.12, MUTED)
    baseline = [("Q projection", "Linear 512→512"), ("K projection", "Linear 512→512"),
                ("V projection", "Linear 512→512"), ("scores", "QKᵀ / √64"),
                ("softmax", "computed in FP32"), ("context", "P·V + output proj")]
    for i, (t, d) in enumerate(baseline):
        x = 0.76 + i * 2.03
        stage_box(slide, x, 2.56, 1.62, t, d, CORAL if i < 3 else AMBER, h=0.94)
        if i < 5: add_arrow(slide, x + 1.62, 3.03, x + 1.92, 3.03, MUTED, 1.5)
    add_label(slide, "Optimized", 0.75, 4.08, 1.35, TEAL)
    stage_box(slide, 1.04, 4.54, 2.38, "Fused QKV projection", "one Linear 512→1536", TEAL, h=1.02)
    add_arrow(slide, 3.42, 5.05, 4.06, 5.05, TEAL)
    stage_box(slide, 4.18, 4.54, 4.05, "scaled_dot_product_attention", "mask + scale + softmax + P·V in a fused backend", CYAN, h=1.02)
    add_arrow(slide, 8.23, 5.05, 8.87, 5.05, CYAN)
    stage_box(slide, 8.99, 4.54, 2.54, "Output projection", "bias optionally deferred", AMBER, h=1.02)
    add_box(slide, 9.25, 6.00, 3.20, 0.70, fill=PANEL_2, line=TEAL)
    add_text(slide, "Why it wins", 9.48, 6.13, 0.95, 0.23, 12, TEAL, True, margin=0)
    add_text(slide, "fewer launches · less HBM traffic", 10.38, 6.12, 1.85, 0.25,
             10.5, INK, margin=0)
    add_footer(slide, 4, "Source: BaselineSelfAttention, OptimizedSelfAttention, README.md")


def add_slide_5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
    add_title(slide, "Residual + bias + mask + LayerNorm: one CUDA pass",
              "A shape-specialized kernel removes intermediate tensors around both residual joins.",
              "03 · Stage optimizations")
    # Before / after.
    add_label(slide, "Native operator chain", 0.75, 2.10, 1.75, CORAL)
    ops = ["projection\nbias", "residual\nadd", "padding\nmask", "LayerNorm\nreduce", "normalized\nwrite"]
    for i, op in enumerate(ops):
        x = 0.78 + i * 1.49
        stage_box(slide, x, 2.58, 1.20, op, "separate op", CORAL, h=0.94)
        if i < 4: add_arrow(slide, x + 1.20, 3.05, x + 1.39, 3.05, MUTED, 1.3)
    add_text(slide, "multiple global-memory round trips", 2.15, 3.68, 4.45, 0.28,
             12, CORAL, True, align=PP_ALIGN.CENTER, margin=0)
    add_arrow(slide, 6.45, 3.08, 7.22, 3.08, TEAL, 2.5)
    add_label(slide, "Custom CUDA", 7.42, 2.10, 1.35, TEAL)
    add_box(slide, 7.42, 2.58, 4.88, 1.38, fill=PANEL, line=TEAL, lw=2)
    add_text(slide, "residual_bias_layer_norm_512", 7.68, 2.82, 4.38, 0.34,
             18, TEAL, True, align=PP_ALIGN.CENTER, margin=0)
    add_text(slide, "residual + delta + bias + mask + mean/variance + affine LN",
             7.72, 3.30, 4.30, 0.27, 11, MUTED, align=PP_ALIGN.CENTER, margin=0)
    # Kernel mechanics cards.
    cards = [
        ("1 block / token", "Each 512-wide row is owned by one 128-thread block.", CYAN),
        ("float4 vector I/O", "Each thread processes four contiguous FP32 values.", TEAL),
        ("warp reductions", "Sum and square-sum stay on-chip for LayerNorm.", AMBER),
        ("two outputs", "Residual state and next normalized input are written once.", CORAL),
    ]
    for i, (t, d, c) in enumerate(cards):
        x = 0.78 + i * 3.03
        add_box(slide, x, 4.62, 2.73, 1.45, fill=PANEL, line=c)
        add_text(slide, t, x + 0.16, 4.84, 2.41, 0.28, 14, c, True, margin=0)
        add_text(slide, d, x + 0.16, 5.27, 2.41, 0.52, 11, MUTED, margin=0)
    add_footer(slide, 5, "Source: cuda_extension/kernels.cu and UserOptimizedTransformer._forward_cuda()")


def add_slide_6(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
    add_title(slide, "FFN: keep elite GEMMs, fuse the activation epilogue",
              "The code accelerates what is memory-bound without replacing cuBLAS/cuBLASLt matmuls.",
              "03 · Stage optimizations")
    stage_box(slide, 0.90, 2.32, 2.18, "Expansion GEMM", "[512] × W₁ → [2048]\ncuBLAS/cuBLASLt", CYAN, h=1.14)
    add_arrow(slide, 3.08, 2.89, 3.72, 2.89, CYAN)
    stage_box(slide, 3.84, 2.32, 2.18, "bias + exact GELU", "custom float4 CUDA kernel\nerf-based exact GELU", TEAL, h=1.14)
    add_arrow(slide, 6.02, 2.89, 6.66, 2.89, TEAL)
    stage_box(slide, 6.78, 2.32, 2.18, "Contraction GEMM", "[2048] × W₂ → [512]\ncuBLAS/cuBLASLt", AMBER, h=1.14)
    add_arrow(slide, 8.96, 2.89, 9.60, 2.89, AMBER)
    stage_box(slide, 9.72, 2.32, 2.66, "Fused residual + LN", "defer output bias into\nnext CUDA fusion", CORAL, h=1.14)
    add_box(slide, 0.90, 4.10, 5.55, 2.02, fill=PANEL, line=GRID)
    add_text(slide, "Why not a handwritten GEMM?", 1.20, 4.40, 4.98, 0.34,
             18, AMBER, True, margin=0)
    add_bullets(slide, [
        "Vendor GEMMs already exploit tensor cores and tuned tiling.",
        "The custom work focuses on bias, activation, masking and normalization.",
    ], 1.20, 4.98, 4.88, 0.82, size=13, bullet_color=AMBER)
    add_box(slide, 6.73, 4.10, 5.65, 2.02, fill=PANEL, line=GRID)
    add_text(slide, "Launch and bandwidth savings", 7.03, 4.40, 5.08, 0.34,
             18, TEAL, True, margin=0)
    add_bullets(slide, [
        "Bias is folded into bias+GELU or residual+LayerNorm.",
        "Fewer temporary allocations and fewer HBM read/write cycles.",
    ], 7.03, 4.98, 4.98, 0.82, size=13)
    add_footer(slide, 6, "Source: README.md, bias_gelu_2048_kernel(), _forward_cuda()")


def add_slide_7(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
    add_title(slide, "Selective TF32: spend error budget where it buys throughput",
              "Hybrid mode uses tensor cores only on the FFN GEMMs that passed the strict accuracy matrix.",
              "04 · Precision strategy")
    add_text(slide, "Layer", 0.80, 2.25, 0.72, 0.28, 11, MUTED, True, margin=0)
    for i in range(6):
        add_text(slide, str(i + 1), 1.72 + i * 1.72, 2.23, 0.35, 0.3, 12,
                 INK, True, align=PP_ALIGN.CENTER, margin=0)
    rows = [
        ("Attention projections", ["FP32"] * 6, CYAN),
        ("FFN expansion W₁", ["TF32", "TF32", "TF32", "FP32", "FP32", "FP32"], TEAL),
        ("FFN contraction W₂", ["TF32"] * 6, AMBER),
    ]
    for r, (label, vals, color) in enumerate(rows):
        y = 2.72 + r * 0.92
        add_text(slide, label, 0.80, y + 0.18, 1.74, 0.28, 11.5, color, True, margin=0)
        for i, val in enumerate(vals):
            fill = PANEL_2 if val == "FP32" else color
            text_color = MUTED if val == "FP32" else BG
            add_box(slide, 2.54 + i * 1.60, y, 1.32, 0.62, fill=fill,
                    line=color if val == "TF32" else GRID)
            add_text(slide, val, 2.54 + i * 1.60, y + 0.17, 1.32, 0.24, 11,
                     text_color, True, align=PP_ALIGN.CENTER, margin=0)
    add_box(slide, 0.80, 5.76, 5.58, 0.76, fill=PANEL, line=CORAL)
    add_text(slide, "Full-model TF32", 1.08, 5.97, 1.58, 0.25, 13, CORAL, True, margin=0)
    add_text(slide, "fails the causal + 35% padding accuracy suite", 2.63, 5.97,
             3.47, 0.25, 11.5, INK, margin=0)
    add_box(slide, 6.68, 5.76, 5.60, 0.76, fill=PANEL, line=TEAL)
    add_text(slide, "Hybrid policy", 6.96, 5.97, 1.32, 0.25, 13, TEAL, True, margin=0)
    add_text(slide, "passes and clears the 3% promotion gate", 8.29, 5.97,
             3.65, 0.25, 11.5, INK, margin=0)
    add_footer(slide, 7, "Source: HYBRID_TF32_FFN_EXPANSION_LAYERS and _forward_cuda(); README.md")


def add_slide_8(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
    add_title(slide, "Correctness gates performance",
              "No backend is promoted on speed alone; every output element must satisfy the explicit OR criterion.",
              "05 · Validation")
    add_box(slide, 0.82, 2.17, 11.68, 1.05, fill=PANEL, line=TEAL, lw=1.8)
    add_rich_text(slide, [
        ("PASS when  ", MUTED, False), ("|user − ref| ≤ 0.001", TEAL, True),
        ("   OR   ", MUTED, False), ("|user − ref| ≤ 0.01 × |ref|", CYAN, True),
    ], 1.17, 2.46, 10.96, 0.34, size=21, align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE)
    steps = [
        ("1", "Identical weights", "Q/K/V weights are concatenated into the optimized QKV matrix.", CYAN),
        ("2", "Five seeded trials", "Random inputs; optional causal and padding-mask stress cases.", TEAL),
        ("3", "Elementwise check", "Counts failures and records worst absolute/relative error.", AMBER),
        ("4", "Benchmark or stop", "Performance is skipped on failure unless explicitly overridden.", CORAL),
    ]
    for i, (num, title, detail, c) in enumerate(steps):
        x = 0.82 + i * 2.98
        add_box(slide, x, 3.78, 2.66, 2.06, fill=PANEL, line=c)
        add_box(slide, x + 0.18, 3.98, 0.48, 0.48, fill=c, line=c)
        add_text(slide, num, x + 0.18, 4.09, 0.48, 0.22, 12, BG, True,
                 align=PP_ALIGN.CENTER, margin=0)
        add_text(slide, title, x + 0.78, 4.04, 1.60, 0.30, 14, c, True, margin=0)
        add_text(slide, detail, x + 0.20, 4.72, 2.26, 0.77, 11.5, MUTED, margin=0)
    add_text(slide, "Validated claim in README: all three reported backends PASS strict accuracy; hybrid also passes causal + 35% padding.",
             1.05, 6.27, 11.05, 0.30, 12.5, INK, bold=True,
             align=PP_ALIGN.CENTER, margin=0)
    add_footer(slide, 8, "Source: compare_outputs(), run_accuracy_tests(), README.md")


def add_slide_9(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
    add_title(slide, "Representative RTX 5090 results",
              "Default FP32 shape; values below are the repository’s documented validation results.",
              "06 · Results")
    # Horizontal latency bars (lower is better), scaled to 1.30 ms.
    add_text(slide, "Median latency ↓", 0.82, 2.14, 2.4, 0.30, 15, INK, True, margin=0)
    results = [
        ("Compiled PyTorch", 1.2036, 1.880, CYAN),
        ("Custom CUDA", 1.1986, 1.889, AMBER),
        ("CUDA hybrid", 0.9323, 2.268, TEAL),
    ]
    for i, (name, latency, speed, color) in enumerate(results):
        y = 2.72 + i * 0.86
        add_text(slide, name, 0.82, y + 0.13, 1.85, 0.27, 12, color, True, margin=0)
        add_box(slide, 2.72, y, 4.35, 0.54, fill=PANEL_2, line=GRID, radius=False)
        bar_w = 4.35 * latency / 1.30
        add_box(slide, 2.72, y, bar_w, 0.54, fill=color, line=color, radius=False)
        add_text(slide, f"{latency:.4f} ms", 2.86, y + 0.13, 1.15, 0.23,
                 11, BG if bar_w > 1.3 else INK, True, margin=0)
    # Result table.
    add_box(slide, 7.48, 2.10, 4.95, 3.05, fill=PANEL, line=GRID)
    headers = ["Backend", "Speedup", "Accuracy"]
    widths = [2.25, 1.12, 1.12]
    x0 = 7.72
    x = x0
    for h, w in zip(headers, widths):
        add_text(slide, h, x, 2.38, w, 0.25, 11, MUTED, True,
                 align=PP_ALIGN.LEFT if h == "Backend" else PP_ALIGN.CENTER, margin=0)
        x += w
    for i, (name, latency, speed, color) in enumerate(results):
        y = 2.92 + i * 0.67
        add_text(slide, name, x0, y, widths[0], 0.27, 12, color, True, margin=0)
        add_text(slide, f"{speed:.3f}×", x0 + widths[0], y, widths[1], 0.27,
                 13, INK, True, align=PP_ALIGN.CENTER, margin=0)
        add_text(slide, "PASS", x0 + widths[0] + widths[1], y, widths[2], 0.27,
                 12, TEAL, True, align=PP_ALIGN.CENTER, margin=0)
    add_box(slide, 0.82, 5.55, 3.55, 0.96, fill=PANEL, line=TEAL)
    add_text(slide, "−22.5%", 1.06, 5.73, 1.32, 0.36, 22, TEAL, True, margin=0)
    add_text(slide, "hybrid latency vs compiled", 2.12, 5.78, 1.95, 0.26,
             10.5, MUTED, margin=0)
    add_box(slide, 4.60, 5.55, 3.55, 0.96, fill=PANEL, line=AMBER)
    add_text(slide, "+0.4%", 4.84, 5.73, 1.30, 0.36, 22, AMBER, True, margin=0)
    add_text(slide, "strict CUDA vs compiled", 5.95, 5.78, 1.90, 0.26,
             10.5, MUTED, margin=0)
    add_box(slide, 8.38, 5.55, 4.05, 0.96, fill=PANEL, line=CYAN)
    add_text(slide, "2.520×", 8.62, 5.73, 1.50, 0.36, 22, CYAN, True, margin=0)
    add_text(slide, "hybrid speedup, causal + 35% padding", 10.18, 5.71, 1.98, 0.38,
             10, MUTED, margin=0)
    add_footer(slide, 9, "Source: README.md representative RTX 5090 validation; lower latency is better")


def add_slide_10(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
    add_title(slide, "What moved the needle—and where it applies",
              "The winning design is a hybrid of mature libraries, narrow custom fusion and measured precision.",
              "07 · Takeaways")
    takeaways = [
        ("Fuse the plumbing", "Residual, bias, padding and LayerNorm are bandwidth-heavy; one pass eliminates temporary traffic.", TEAL),
        ("Trust tuned primitives", "Fused SDPA and cuBLAS/cuBLASLt remain the right engines for attention and GEMMs.", CYAN),
        ("Use precision surgically", "Selective TF32 captures tensor-core gains while keeping sensitive attention paths in strict FP32.", AMBER),
    ]
    for i, (title, detail, color) in enumerate(takeaways):
        x = 0.78 + i * 4.06
        add_box(slide, x, 2.23, 3.72, 2.10, fill=PANEL, line=color, lw=1.5)
        add_text(slide, f"0{i+1}", x + 0.22, 2.50, 0.48, 0.34, 18, color, True, margin=0)
        add_text(slide, title, x + 0.82, 2.50, 2.62, 0.34, 17, color, True, margin=0)
        add_text(slide, detail, x + 0.24, 3.10, 3.22, 0.78, 12.5, MUTED, margin=0)
    add_box(slide, 0.78, 4.76, 11.84, 1.45, fill=PANEL_2, line=CORAL)
    add_text(slide, "Scope boundary", 1.06, 5.04, 1.60, 0.30, 16, CORAL, True, margin=0)
    add_text(slide, "The custom CUDA backend is inference-only and specialized for CUDA FP32 at 8×128×512, 8 heads, 2048 FFN and 6 layers. Other shapes/dtypes fall back to the portable PyTorch path.",
             2.55, 4.98, 9.65, 0.66, 13, INK, margin=0)
    add_text(slide, "Bottom line: 0.9323 ms median · 2.268× baseline speedup · strict accuracy PASS",
             1.05, 6.52, 11.20, 0.32, 15, TEAL, True, align=PP_ALIGN.CENTER, margin=0)
    add_footer(slide, 10, "Source: README.md and custom_cuda_shape_supported()")


def add_slide_11(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
    add_label(slide, "Appendix", 0.72, 0.64, 1.10, CYAN)
    add_text(slide, "Reproduce the benchmark", 0.72, 1.20, 7.5, 0.62, 30, INK, True, margin=0)
    add_text(slide, "Build the extension, compare all backends, then run the padded causal stress case.",
             0.72, 1.92, 9.6, 0.36, 14, MUTED, margin=0)
    commands = [
        ("1", "make build-cuda"),
        ("2", "make benchmark-all"),
        ("3", "make benchmark-best"),
    ]
    for i, (n, cmd) in enumerate(commands):
        y = 2.72 + i * 1.06
        add_box(slide, 0.74, y, 0.58, 0.58, fill=TEAL, line=TEAL)
        add_text(slide, n, 0.74, y + 0.15, 0.58, 0.22, 13, BG, True,
                 align=PP_ALIGN.CENTER, margin=0)
        add_box(slide, 1.58, y, 7.55, 0.58, fill=PANEL, line=GRID)
        add_text(slide, cmd, 1.82, y + 0.14, 7.05, 0.26, 15, INK, False,
                 font="DejaVu Sans Mono", margin=0)
    add_box(slide, 9.53, 2.72, 2.92, 2.70, fill=PANEL, line=CYAN)
    add_text(slide, "Reported output", 9.82, 3.02, 2.34, 0.31, 17, CYAN, True, margin=0)
    add_bullets(slide, [
        "Accuracy PASS/FAIL",
        "median / mean / p90 / min",
        "token throughput",
        "baseline speedup",
        "peak activation allocation",
    ], 9.82, 3.58, 2.30, 1.42, size=11.5, bullet_color=CYAN, gap=4)
    add_text(slide, "Generated from the repository implementation and its documented representative results.",
             0.74, 6.64, 10.5, 0.26, 10.5, MUTED, margin=0)
    add_text(slide, "11", 12.05, 6.63, 0.58, 0.28, 10, MUTED, True,
             align=PP_ALIGN.RIGHT, margin=0)


def build_deck():
    prs = Presentation()
    prs.slide_width = SW; prs.slide_height = SH
    prs.core_properties.title = "Transformer Inference Optimization Pipeline and Results"
    prs.core_properties.subject = "PyTorch Transformer GPU optimization on RTX 5090"
    prs.core_properties.author = "TikTok TechJam 2026"
    prs.core_properties.keywords = "Transformer, CUDA, SDPA, TF32, GPU optimization"
    for builder in [add_slide_1, add_slide_2, add_slide_3, add_slide_4,
                    add_slide_5, add_slide_6, add_slide_7, add_slide_8,
                    add_slide_9, add_slide_10, add_slide_11]:
        builder(prs)
    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    print(build_deck())
